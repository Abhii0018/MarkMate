import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    # Color Palette: Deep Navy, Modern Teal, Accent Coral, Card Dark, White
    COLOR_BG = RGBColor(15, 23, 42)        # #0F172A (Slate 900)
    COLOR_CARD = RGBColor(30, 41, 59)      # #1E293B (Slate 800)
    COLOR_PRIMARY = RGBColor(59, 130, 246)  # #3B82F6 (Primary Blue)
    COLOR_ACCENT = RGBColor(239, 68, 68)    # #EF4444 (Red Accent / Threshold Alert)
    COLOR_TEXT = RGBColor(241, 245, 249)    # #F1F5F9 (Light Grey/White)
    COLOR_MUTED = RGBColor(148, 163, 184)   # #94A3B8 (Slate Muted)
    COLOR_HIGHLIGHT = RGBColor(16, 185, 129)# #10B981 (Emerald Green)
    COLOR_INDIGO = RGBColor(99, 102, 241)   # #6366F1 (Indigo / Teacher Theme)

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_text="MARKMATE MULTI-PORTAL ATTENDANCE SYSTEM"):
        # Category label
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_PRIMARY
        p_cat.font.name = "Calibri"

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT
        p_title.font.name = "Calibri"

    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD, border_color=COLOR_CARD):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        return shape

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1, COLOR_BG)

    # Hero card container
    add_card(slide1, Inches(1.0), Inches(1.1), Inches(11.333), Inches(5.3), bg_color=COLOR_CARD)
    
    tbox = slide1.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.3), Inches(2.2))
    tf = tbox.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "MARKMATE"
    p0.font.size = Pt(44)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_PRIMARY
    p0.font.name = "Calibri"
    
    p1 = tf.add_paragraph()
    p1.text = "Multi-Portal Student Attendance & Department Management System"
    p1.font.size = Pt(24)
    p1.font.color.rgb = COLOR_TEXT
    p1.font.bold = True
    p1.font.name = "Calibri"
    
    p2 = tf.add_paragraph()
    p2.text = "An automated, full-stack institution portal supporting Civil, Mechanical, Computer Science departments with dedicated Admin, Teacher, and Student portals and automated <75% RED threshold alerts."
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_MUTED
    p2.space_before = Pt(12)
    p2.font.name = "Calibri"

    # Presenter Details & Highlights
    tbox_info = slide1.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(10.3), Inches(1.6))
    tf_info = tbox_info.text_frame
    
    p_info1 = tf_info.paragraphs[0]
    p_info1.text = "✨ Portals: Admin Dashboard • Teacher Portal (Dept-based) • Student Portal (<75% RED Alert)"
    p_info1.font.size = Pt(14)
    p_info1.font.bold = True
    p_info1.font.color.rgb = COLOR_HIGHLIGHT
    p_info1.font.name = "Calibri"

    p_info2 = tf_info.add_paragraph()
    p_info2.text = "🛠️ Tech Stack: Python Flask • SQLAlchemy ORM • Dual MySQL/SQLite • Bootstrap 5 Glassmorphism"
    p_info2.font.size = Pt(13)
    p_info2.font.color.rgb = COLOR_MUTED
    p_info2.space_before = Pt(6)
    p_info2.font.name = "Calibri"

    # ==========================================
    # SLIDE 2: Problem Statement & Solution
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2, COLOR_BG)
    add_header(slide2, "Problem Statement & Solution")

    # Left Card - The Problem
    add_card(slide2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    tb_p = slide2.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.4))
    tf_p = tb_p.text_frame
    tf_p.word_wrap = True
    
    p = tf_p.paragraphs[0]
    p.text = "❌ Traditional Academic Challenges"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    
    points_p = [
        "Time-Consuming Roll Calls: Manual registers take away 10-15 minutes of active lecture time.",
        "Lack of Department Isolation: Faculty teachers cannot filter or manage students by Civil, Mechanical, or Computer Science.",
        "Unaware of Attendance Shortages: Students discover attendance shortages (<75%) too late at exam time.",
        "No Centralized Student Visibility: Lack of self-service portals for students to inspect daily attendance history."
    ]
    for pt in points_p:
        p_sub = tf_p.add_paragraph()
        p_sub.text = "• " + pt
        p_sub.font.size = Pt(13)
        p_sub.font.color.rgb = COLOR_TEXT
        p_sub.space_before = Pt(10)

    # Right Card - The MarkMate Solution
    add_card(slide2, Inches(6.9), Inches(1.8), Inches(5.6), Inches(5.0))
    tb_s = slide2.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.0), Inches(4.4))
    tf_s = tb_s.text_frame
    tf_s.word_wrap = True
    
    p = tf_s.paragraphs[0]
    p.text = "💡 The MarkMate Solution"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_HIGHLIGHT
    
    points_s = [
        "Multi-Role Portals: Role-based access for Admin, Department Teachers, and Students.",
        "Department-Based Marking: Teachers log in and mark attendance specifically for Civil, ME, CS students.",
        "Automated <75% RED Threshold Alert: Instant red alert card & gauge highlight when student attendance <75%.",
        "Live Dashboard Analytics: Real-time counter metrics, present/absent stats, and exportable logs."
    ]
    for pt in points_s:
        p_sub = tf_s.add_paragraph()
        p_sub.text = "• " + pt
        p_sub.font.size = Pt(13)
        p_sub.font.color.rgb = COLOR_TEXT
        p_sub.space_before = Pt(10)

    # ==========================================
    # SLIDE 3: Three Dedicated Portals Architecture
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3, COLOR_BG)
    add_header(slide3, "Multi-Portal Ecosystem & Access Control")

    portals = [
        ("🛡️ Admin Dashboard", "Centralized System Management", [
            "• Manage faculty teachers (Add/Edit/Delete)",
            "• Filter students per department (Civil, ME, CS, IT)",
            "• Department breakdown cards & teacher counts",
            "• System-wide daily attendance overview"
        ], COLOR_PRIMARY),
        ("👨‍🏫 Teacher Portal", "Department-Isolated Interface", [
            "• Dedicated login for Civil, ME, CS teachers",
            "• Mark daily attendance for department students",
            "• Department attendance rate analytics",
            "• Historical logs filterable by date"
        ], COLOR_INDIGO),
        ("🎓 Student Portal", "Self-Service & Threshold Alert", [
            "• Individual roll number authentication",
            "• Overall percentage calculation",
            "• RED Alert Banner if attendance < 75%",
            "• Full daily attendance record history"
        ], COLOR_ACCENT)
    ]

    for idx, (p_name, p_sub, p_list, p_color) in enumerate(portals):
        x = Inches(0.8 + idx * 4.0)
        add_card(slide3, x, Inches(1.8), Inches(3.7), Inches(5.0))
        
        tb = slide3.shapes.add_textbox(x + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = p_name
        p0.font.size = Pt(18)
        p0.font.bold = True
        p0.font.color.rgb = p_color
        
        p1 = tf.add_paragraph()
        p1.text = p_sub
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_MUTED
        p1.space_before = Pt(4)

        for item in p_list:
            pf = tf.add_paragraph()
            pf.text = item
            pf.font.size = Pt(12)
            pf.font.color.rgb = COLOR_TEXT
            pf.space_before = Pt(10)

    # ==========================================
    # SLIDE 4: Department Faculty & Student Management
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4, COLOR_BG)
    add_header(slide4, "Department-Wise Teacher & Student Management")

    depts = [
        ("💻 Computer Science", "Faculty: Prof. Rajesh Sharma & Prof. Vikram Singh", "Manages CS101, CS102, CS103, CS104 student rosters & semester attendance."),
        ("⚙️ Mechanical Engineering", "Faculty: Prof. Suresh Verma", "Manages ME101, ME102 student rosters and laboratory attendance records."),
        ("🏗️ Civil Engineering", "Faculty: Prof. Ananya Gupta", "Manages CE101, CE102 student rosters and structural batch attendance."),
        ("🌐 Information Technology", "Faculty: Department Assigned", "Manages IT101, IT102, IT103 student rosters and practical sessions.")
    ]

    for idx, (d_name, d_fac, d_desc) in enumerate(depts):
        x = Inches(0.8 if idx % 2 == 0 else 6.9)
        y = Inches(1.8 if idx < 2 else 4.5)
        
        add_card(slide4, x, y, Inches(5.6), Inches(2.4))
        tb = slide4.shapes.add_textbox(x + Inches(0.3), y + Inches(0.25), Inches(5.0), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = d_name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
        p_f = tf.add_paragraph()
        p_f.text = d_fac
        p_f.font.size = Pt(13)
        p_f.font.bold = True
        p_f.font.color.rgb = COLOR_HIGHLIGHT
        p_f.space_before = Pt(6)

        p_d = tf.add_paragraph()
        p_d.text = d_desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = COLOR_MUTED
        p_d.space_before = Pt(6)

    # ==========================================
    # SLIDE 5: 75% Attendance Red Threshold Engine
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5, COLOR_BG)
    add_header(slide5, "Automated 75% Attendance Red Threshold Alert")

    # Left Column - Calculation Formula & Rule
    add_card(slide5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    tb_l = slide5.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.4))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True

    p = tf_l.paragraphs[0]
    p.text = "📐 Threshold Calculation Engine"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    p_f = tf_l.add_paragraph()
    p_f.text = "Formula: Percentage = (Attended Classes / Total Conducted) * 100"
    p_f.font.size = Pt(13)
    p_f.font.bold = True
    p_f.font.color.rgb = COLOR_HIGHLIGHT
    p_f.space_before = Pt(12)

    rules = [
        "Threshold Rule: If overall attendance < 75.0%, trigger RED Alert UI.",
        "Compliant Rule: If overall attendance >= 75.0%, display Green Good Standing.",
        "Automatic Evaluation: Calculated dynamically on every student dashboard access.",
        "Early Intervention: Enables students to take corrective action prior to exam eligibility cutoff."
    ]
    for r in rules:
        pr = tf_l.add_paragraph()
        pr.text = "• " + r
        pr.font.size = Pt(12)
        pr.font.color.rgb = COLOR_TEXT
        pr.space_before = Pt(10)

    # Right Column - Visual UI Reactions
    add_card(slide5, Inches(6.9), Inches(1.8), Inches(5.6), Inches(5.0), bg_color=COLOR_CARD, border_color=COLOR_ACCENT)
    tb_r = slide5.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.0), Inches(4.4))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "⚠️ Red Threshold Visual Reaction (< 75%)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    reactions = [
        "1. Pulsing Red Alert Banner: Prominent warning box notifying the exact shortage.",
        "2. Red Stat Gauge: Stat percentage card highlights in vivid Red (e.g. 33.3%).",
        "3. Red Progress Bar: High-contrast red bar indicating critical threshold shortage.",
        "4. Badge Highlight: 'Below 75% Threshold' badge displayed on profile header."
    ]
    for rx in reactions:
        prx = tf_r.add_paragraph()
        prx.text = rx
        prx.font.size = Pt(13)
        prx.font.color.rgb = COLOR_TEXT
        prx.space_before = Pt(12)

    # ==========================================
    # SLIDE 6: Technology Stack & Security
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6, COLOR_BG)
    add_header(slide6, "Technology Architecture & Security Stack")

    tech_stack = [
        ("Backend Framework", "Python 3 & Flask Framework", "Flask Blueprints (Auth, Admin, Teachers, Teacher Portal, Student Portal, Attendance).", COLOR_PRIMARY),
        ("Database & ORM", "SQLAlchemy & Dual MySQL/SQLite", "SQLAlchemy ORM with relational foreign key cascading and date uniqueness constraints.", COLOR_HIGHLIGHT),
        ("Security & Auth", "Werkzeug Hashing & Role Decorators", "Scrypt password hashing, session management, `@admin_required`, `@teacher_required`, `@student_required`.", COLOR_ACCENT),
        ("Frontend & Styling", "HTML5, Vanilla CSS & Bootstrap 5", "Modern glassmorphism UI, Outfit/Plus Jakarta typography, animated counter numbers.", COLOR_INDIGO)
    ]

    for idx, (layer, title, desc, color) in enumerate(tech_stack):
        y = Inches(1.8 + idx * 1.3)
        add_card(slide6, Inches(0.8), y, Inches(11.7), Inches(1.1))
        
        tb_left = slide6.shapes.add_textbox(Inches(1.1), y + Inches(0.15), Inches(3.2), Inches(0.8))
        tf_l = tb_left.text_frame
        tf_l.word_wrap = True
        p_l = tf_l.paragraphs[0]
        p_l.text = layer.upper()
        p_l.font.size = Pt(11)
        p_l.font.bold = True
        p_l.font.color.rgb = color
        
        p_l2 = tf_l.add_paragraph()
        p_l2.text = title
        p_l2.font.size = Pt(15)
        p_l2.font.bold = True
        p_l2.font.color.rgb = COLOR_TEXT

        tb_right = slide6.shapes.add_textbox(Inches(4.5), y + Inches(0.25), Inches(7.7), Inches(0.7))
        tf_r = tb_right.text_frame
        tf_r.word_wrap = True
        p_r = tf_r.paragraphs[0]
        p_r.text = desc
        p_r.font.size = Pt(13)
        p_r.font.color.rgb = COLOR_MUTED

    # ==========================================
    # SLIDE 7: Database Schema & Entity Models
    # ==========================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide7, COLOR_BG)
    add_header(slide7, "Database Schema & Relational Models")

    models = [
        ("Teacher Entity", "Table: `teachers`", ["• id (PK, Integer)", "• username (Unique)", "• password (Hashed)", "• name (String)", "• department (String)", "• email (String)"], COLOR_INDIGO),
        ("Student Entity", "Table: `students`", ["• id (PK, Integer)", "• roll_number (Unique)", "• name (String)", "• department (String)", "• semester (String)", "• password (Hashed)"], COLOR_HIGHLIGHT),
        ("Attendance Entity", "Table: `attendance`", ["• id (PK, Integer)", "• student_id (FK)", "• attendance_date (Date)", "• status ('Present'/'Absent')", "• Unique(student_id, date)"], COLOR_ACCENT)
    ]

    for idx, (m_name, m_table, m_fields, m_color) in enumerate(models):
        x = Inches(0.8 + idx * 4.0)
        add_card(slide7, x, Inches(1.8), Inches(3.7), Inches(5.0))
        
        tb = slide7.shapes.add_textbox(x + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = m_name
        p0.font.size = Pt(18)
        p0.font.bold = True
        p0.font.color.rgb = m_color
        
        p1 = tf.add_paragraph()
        p1.text = m_table
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_MUTED
        p1.space_before = Pt(4)
        
        p_space = tf.add_paragraph()
        p_space.text = "Attributes:"
        p_space.font.size = Pt(13)
        p_space.font.bold = True
        p_space.font.color.rgb = COLOR_TEXT
        p_space.space_before = Pt(12)

        for field in m_fields:
            pf = tf.add_paragraph()
            pf.text = field
            pf.font.size = Pt(12)
            pf.font.color.rgb = COLOR_TEXT
            pf.space_before = Pt(6)

    # ==========================================
    # SLIDE 8: Future Enhancements & Roadmap
    # ==========================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide8, COLOR_BG)
    add_header(slide8, "Future Roadmap & Project Scaling")

    enhancements = [
        ("📲 Automated SMS / WhatsApp Alerts", "Instant automated WhatsApp alerts to parents when student attendance drops below 75%.", COLOR_PRIMARY),
        ("📸 AI Biometric / Facial Recognition", "Zero-touch classroom attendance marking via camera face recognition model.", COLOR_ACCENT),
        ("📊 Course & Subject-Wise Tracking", "Granular attendance tracking per individual course subject code and lecture time.", COLOR_HIGHLIGHT),
        ("📱 Native Mobile Application", "iOS & Android mobile apps with push notifications for instant attendance updates.", COLOR_INDIGO)
    ]

    for idx, (title, desc, color) in enumerate(enhancements):
        x = Inches(0.8 if idx % 2 == 0 else 6.9)
        y = Inches(1.8 if idx < 2 else 4.5)
        
        add_card(slide8, x, y, Inches(5.6), Inches(2.4))
        tb = slide8.shapes.add_textbox(x + Inches(0.3), y + Inches(0.3), Inches(5.0), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_TEXT
        p2.space_before = Pt(10)

    # ==========================================
    # SLIDE 9: Conclusion & Q&A Slide
    # ==========================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide9, COLOR_BG)

    add_card(slide9, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5), bg_color=COLOR_CARD)
    
    tb_c = slide9.shapes.add_textbox(Inches(2.0), Inches(2.0), Inches(9.333), Inches(3.5))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    
    pc1 = tf_c.paragraphs[0]
    pc1.alignment = PP_ALIGN.CENTER
    pc1.text = "Thank You!"
    pc1.font.size = Pt(42)
    pc1.font.bold = True
    pc1.font.color.rgb = COLOR_PRIMARY
    
    pc2 = tf_c.add_paragraph()
    pc2.alignment = PP_ALIGN.CENTER
    pc2.text = "MarkMate • Multi-Portal Student Attendance & Department System"
    pc2.font.size = Pt(20)
    pc2.font.bold = True
    pc2.font.color.rgb = COLOR_TEXT
    pc2.space_before = Pt(14)

    pc3 = tf_c.add_paragraph()
    pc3.alignment = PP_ALIGN.CENTER
    pc3.text = "Questions & Answers (Q&A)"
    pc3.font.size = Pt(24)
    pc3.font.bold = True
    pc3.font.color.rgb = COLOR_HIGHLIGHT
    pc3.space_before = Pt(24)

    # Save presentation
    output_path = "/Users/abhishekkumar/Desktop/Attendance/MarkMate_Project_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation successfully created at: {output_path}")

if __name__ == "__main__":
    create_presentation()
